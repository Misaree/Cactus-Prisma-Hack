!pip install python-pptx
!pip install pdf2image
!pip install pymupdf

from pptx import Presentation
from pptx.util import Inches, Pt
import spacy
import os
import fitz  # PyMuPDF
from PIL import Image, ImageDraw, ImageFont

def generate_presentation_with_images(pdf_file, summarized_text, output_folder, output_file):
    # Initialize spaCy model for NLP
    nlp = spacy.load("en_core_web_sm")

    # Helper function to generate title from the filename
    def generate_title_from_filename(filename):
        base_name = os.path.splitext(os.path.basename(filename))[0]  # Remove extension
        doc = nlp(base_name.replace("_", " "))  # Process text with spaCy
        keywords = [chunk.text for chunk in doc.noun_chunks]
        title = " ".join(keywords[:4]) if keywords else base_name[:4]  # Default to filename if no keywords
        return title.title()  # Capitalize each word

    # Helper function to classify text into sections based on keywords
    def classify_text(text, keywords):
        classified_data = {section: [] for section in keywords}
        doc = nlp(text)
        for sentence in doc.sents:
            sentence_text = sentence.text.strip()
            for section, words in keywords.items():
                if any(word in sentence_text.lower() for word in words):
                    classified_data[section].append(sentence_text)
                    break
        return classified_data

    # Helper function to adjust font size dynamically based on text length
    def get_font_size(text_length):
        if text_length <= 200:
            return Pt(28)  # Large font for short text
        elif text_length <= 400:
            return Pt(26)  # Medium font
        elif text_length <= 600:
            return Pt(24)  # Smaller font
        elif text_length <= 800:
            return Pt(20)  # Even smaller
        else:
            return Pt(16)  # Smallest font for long content

    # Extract images from the PDF and add captions to them
    def extract_and_caption_images(pdf_path, output_folder):
        doc = fitz.open(pdf_path)
        os.makedirs(output_folder, exist_ok=True)

        def add_caption_to_image(image_path, caption_text, output_path):
            with Image.open(image_path) as img:
                width, height = img.size
                caption_height = 50  # Space for the caption
                new_image = Image.new("RGB", (width, height + caption_height), "white")
                new_image.paste(img, (0, 0))
                draw = ImageDraw.Draw(new_image)
                font = ImageFont.load_default()
                draw.text((10, height + 10), caption_text, fill="black", font=font)
                new_image.save(output_path)

        for page_num in range(len(doc)):
            page = doc[page_num]
            images = page.get_images(full=True)
            for img_index, img in enumerate(images):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                image_filename = os.path.join(output_folder, f"page_{page_num + 1}img{img_index + 1}.{image_ext}")
                with open(image_filename, "wb") as f:
                    f.write(image_bytes)
                caption_text = f"Page {page_num + 1}, Image {img_index + 1}"
                img_rect = fitz.Rect(img[1:5])
                page_text = page.get_text("dict")
                for block in page_text["blocks"]:
                    block_rect = fitz.Rect(block["bbox"])
                    if img_rect.intersects(block_rect):
                        for line in block["lines"]:
                            for span in line["spans"]:
                                caption_text += f" {span['text']}"
                captioned_image_filename = image_filename.replace(f".{image_ext}", "_captioned.png")
                add_caption_to_image(image_filename, caption_text.strip(), captioned_image_filename)
                os.remove(image_filename)

        print("Image extraction and captioning completed.")

    # Define keywords for each section
    section_keywords = {
        "Abstract": ["abstract", "summary"],
        "Introduction": ["introduction", "potential", "overview", "functionality"],
        "Literature Review": ["literature review", "related work"],
        "Methodology": ["methodology", "methods", "approach"],
        "Challenges": ["challenges", "accuracy", "compliance", "standards"],
        "SE Assistant Tool": ["SE Assistant", "M-RAG", "tool", "DesDocs"],
        "Evaluation Process": ["evaluation", "criteria", "engineers", "validation"],
        "Future Work": ["future", "enhancements", "next steps", "improvements"],
        "Results": ["results", "findings", "observations"],
        "Conclusion": ["conclusion", "future work", "summary"]
    }

    # Classify the summary into sections
    classified_sections = classify_text(summarized_text, section_keywords)

    # Create the PowerPoint presentation
    prs = Presentation()

    # Add title slide
    presentation_title = generate_title_from_filename(pdf_file)
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = presentation_title
    slide.placeholders[1].text = "Automated Presentation from Research Summary"

    # Add slides dynamically for each section
    for section, content in classified_sections.items():
        if content:
            slide_layout = prs.slide_layouts[1]  # Use content slide layout
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = section
            text_box = slide.shapes.placeholders[1]
            text_frame = text_box.text_frame
            text_frame.clear()  # Clear default text
            total_text = " ".join(content)
            font_size = get_font_size(len(total_text))
            for sentence in content:
                paragraph = text_frame.add_paragraph()
                paragraph.text = sentence
                paragraph.font.size = font_size
                paragraph.space_after = Pt(6)

    # Add "Thank You" slide
    thank_you_slide_layout = prs.slide_layouts[6]  # Blank Slide Layout
    thank_you_slide = prs.slides.add_slide(thank_you_slide_layout)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    left = (slide_width - Inches(3)) / 2
    top = (slide_height - Inches(1.5)) / 2
    width = Inches(3)
    height = Inches(1.5)
    thank_you_textbox = thank_you_slide.shapes.add_textbox(left, top, width, height)
    text_frame = thank_you_textbox.text_frame
    text_frame.text = "Thank You!"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.alignment = 1  # Center align
    footer_width = Inches(3)
    footer_height = Inches(0.5)
    footer_left = slide_width - footer_width - Inches(0.5)
    footer_top = slide_height - footer_height - Inches(0.3)
    footer = thank_you_slide.shapes.add_textbox(footer_left, footer_top, footer_width, footer_height)
    footer_text_frame = footer.text_frame
    footer_text_frame.text = "This PPT was generated by AI"
    p = footer_text_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.bold = True
    p.alignment = 2  # Right align

    # Extract and caption images from the PDF file
    extract_and_caption_images(pdf_file, output_folder)

    # Get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Get all extracted images
    image_files = sorted(os.listdir(output_folder))
    image_files = [os.path.join(output_folder, img) for img in image_files if img.endswith((".png", ".jpg", ".jpeg"))]

    # Insert each image as a new slide before "Thank You"
    for image in image_files:
        with Image.open(image) as img:
            img_width, img_height = img.size
        img_width_pts = Inches(img_width / 96)
        img_height_pts = Inches(img_height / 96)
        max_width = slide_width * 0.7
        max_height = slide_height * 0.7
        if img_width_pts > max_width or img_height_pts > max_height:
            scale_ratio = min(max_width / img_width_pts, max_height / img_height_pts)
            img_width_pts *= scale_ratio
            img_height_pts *= scale_ratio
        left = (slide_width - img_width_pts) / 2
        top = (slide_height - img_height_pts) / 2
        slide_layout = prs.slide_layouts[5]  # Blank slide
        slide = prs.slides.add_slide(slide_layout)
        prs.slides._sldIdLst.insert(-1, prs.slides._sldIdLst[-1])
        slide.shapes.add_picture(image, left, top, width=img_width_pts, height=img_height_pts)

    # Save the presentation
    prs.save(output_file)
    print(f"Presentation created and saved as '{output_file}'!")


from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import spacy
import os

def generate_informal_ppt(summary, pdf_file):
    # Load spaCy model for NLP
    nlp = spacy.load("en_core_web_sm")

    def generate_title_from_filename(filename):
        base_name = os.path.splitext(os.path.basename(filename))[0]
        doc = nlp(base_name.replace("_", " "))
        keywords = [chunk.text for chunk in doc.noun_chunks]
        return " ".join(keywords[:4]).title() if keywords else base_name[:4].title()

    section_keywords = {
        "Abstract": ["abstract", "summary"],
        "Introduction": ["introduction", "potential", "overview", "functionality"],
        "Literature Review": ["literature review", "related work"],
        "Methodology": ["methodology", "methods", "approach"],
        "Challenges": ["challenges", "accuracy", "compliance", "standards"],
        "SE Assistant Tool": ["SE Assistant", "M-RAG", "tool", "DesDocs"],
        "Evaluation Process": ["evaluation", "criteria", "engineers", "validation"],
        "Future Work": ["future", "enhancements", "next steps", "improvements"],
        "Results": ["results", "findings", "observations"],
        "Conclusion": ["conclusion", "future work", "summary"]
    }

    def classify_text(text, keywords):
        classified_data = {section: [] for section in keywords}
        doc = nlp(text)
        for sentence in doc.sents:
            for section, words in keywords.items():
                if any(word in sentence.text.lower() for word in words):
                    classified_data[section].append(sentence.text.strip())
                    break
        return classified_data

    def get_font_size(text_length):
        if text_length <= 200:
            return Pt(28)
        elif text_length <= 400:
            return Pt(26)
        elif text_length <= 600:
            return Pt(24)
        elif text_length <= 800:
            return Pt(20)
        else:
            return Pt(16)

    classified_sections = classify_text(summary, section_keywords)
    prs = Presentation()

    # Title Slide
    presentation_title = generate_title_from_filename(pdf_file)
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = presentation_title
    slide.placeholders[1].text = "Automated Presentation from Research Summary"
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0, 51, 102)
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(44)
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Content Slides
    for section, content in classified_sections.items():
        if content:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = section
            text_box = slide.shapes.placeholders[1]
            text_frame = text_box.text_frame
            text_frame.clear()
            total_text = " ".join(content)
            font_size = get_font_size(len(total_text))

            for sentence in content:
                paragraph = text_frame.add_paragraph()
                paragraph.text = sentence
                paragraph.font.size = font_size
                paragraph.space_after = Pt(6)

            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(204, 229, 255)

    # Thank You Slide
    thank_you_slide_layout = prs.slide_layouts[6]
    thank_you_slide = prs.slides.add_slide(thank_you_slide_layout)
    thank_you_slide.background.fill.solid()
    thank_you_slide.background.fill.fore_color.rgb = RGBColor(0, 51, 102)

    slide_width, slide_height = prs.slide_width, prs.slide_height
    left, top, width, height = (slide_width - Inches(3)) / 2, (slide_height - Inches(1.5)) / 2, Inches(3), Inches(1.5)

    thank_you_textbox = thank_you_slide.shapes.add_textbox(left, top, width, height)
    text_frame = thank_you_textbox.text_frame
    text_frame.text = "Thank You!"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    footer = thank_you_slide.shapes.add_textbox(slide_width - Inches(3.5), slide_height - Inches(0.8), Inches(3), Inches(0.5))
    footer_text_frame = footer.text_frame
    footer_text_frame.text = "This PPT was generated by AI"
    p = footer_text_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(128, 128, 128)
    p.bold = True
    p.alignment = PP_ALIGN.RIGHT

    output_file = "Generative_AI_Presentation_Informal.pptx"
    prs.save(output_file)
    print(f"Presentation created and saved as '{output_file}'!")
    return output_file


choice = input("Do you want to generate a formal or informal PPT? (Enter 'formal' or 'informal'): ").strip().lower()
pdf_file= pdf_file
summarized_text= summarized_text
output_file="Generative_AI_Presentation.pptx"

# Call the respective function
if choice == "formal":
    generate_presentation_with_images(pdf_file, summarized_text, "/content/extracted_images_with_text_final", "Generative_AI_Presentation.pptx")
elif choice == "informal":
    generate_informal_ppt(summarized_text,pdf_file)
else:
    print("Invalid choice! Please enter 'formal' or 'informal'.")
